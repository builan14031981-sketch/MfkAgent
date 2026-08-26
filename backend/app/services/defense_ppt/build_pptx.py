"""模块4：PPT 生成与排版 — 用 python-pptx 把结构化内容渲染为 .pptx。

设计系统（取自平台 UI 设计类 Skill 的核心原则）：
- 中性底色 + 单一克制强调色（accent），不使用 emoji；
- 清晰字号层级：封面/章节 > 页标题 > 正文；
- 充足留白、统一栅格、页脚标注来源；
- 封面/章节页若有真实 AI 图，则满版铺底 + 暗化遮罩保证文字可读；
- 内容页右侧插图来自 assets/content*.png（程序化图表优先于 AI 图表现数据）。
"""
from __future__ import annotations

import json
import os
import re
import shutil
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

_HERE = os.path.dirname(os.path.abspath(__file__))
_STYLES_PATH = os.path.join(_HERE, "templates", "styles.json")
_REAL_DIR = os.path.join(_HERE, "templates", "real")

IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".bmp", ".gif")

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK = RGBColor(0x1A, 0x1A, 0x1A)


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.replace("#", ""))


def _set_font(run, name: str, size: int, color: RGBColor, bold: bool = False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = rPr.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea", {}
        )
        rPr.append(ea)
    ea.set("{http://schemas.openxmlformats.org/drawingml/2006/main}typeface", name)


def _load_styles() -> Dict:
    with open(_STYLES_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def _list_images(assets_dir: Optional[str]) -> List[str]:
    if not assets_dir or not os.path.isdir(assets_dir):
        return []
    imgs = []
    for fn in sorted(os.listdir(assets_dir)):
        if fn.lower().endswith(IMAGE_EXTS):
            imgs.append(os.path.join(assets_dir, fn))
    return imgs


def _named_image(assets_dir: Optional[str], base: str) -> Optional[str]:
    if not assets_dir or not os.path.isdir(assets_dir):
        return None
    for fn in sorted(os.listdir(assets_dir)):
        name, ext = os.path.splitext(fn)
        if name.lower() == base.lower() and ext.lower() in IMAGE_EXTS:
            return os.path.join(assets_dir, fn)
    return None


def _resolve_images(assets_dir: Optional[str]) -> Dict:
    """按契约解析素材图：cover.* / section.* / content*.png。"""
    if not assets_dir or not os.path.isdir(assets_dir):
        return {"cover": None, "section": None, "content": []}
    all_imgs = _list_images(assets_dir)
    cover = _named_image(assets_dir, "cover")
    section = _named_image(assets_dir, "section")
    contents = [
        os.path.join(assets_dir, fn)
        for fn in sorted(os.listdir(assets_dir))
        if fn.lower().startswith("content") and fn.lower().endswith(IMAGE_EXTS)
    ]
    # 兜底：没有按命名放图时，用目录里的图按顺序顶替，保证测试可用
    if not cover and len(all_imgs) >= 1:
        cover = all_imgs[0]
    if not section and len(all_imgs) >= 2:
        section = all_imgs[1]
    if not contents:
        rest = [i for i in all_imgs if i not in (cover, section)]
        contents = rest
    return {"cover": cover, "section": section, "content": contents}


def _real_master(discipline: str, style: str) -> Optional[str]:
    cand = os.path.join(_REAL_DIR, f"{discipline}_{style}.pptx")
    return cand if os.path.exists(cand) else None


# ───────────────────────── 基础图元 ─────────────────────────

def _bg(slide, hex_color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _rect(slide, l, t, w, h, hex_color, line_hex=None, line_w=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(hex_color)
    if line_hex is not None:
        shp.line.color.rgb = _rgb(line_hex)
        shp.line.width = Pt(line_w or 1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _rect_alpha(slide, l, t, w, h, hex_color, alpha: float):
    """半透明矩形（用于图片遮罩）。alpha: 0~1。"""
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(hex_color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    sf = shp.fill._xPr.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None:
            a = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
            srgb.append(a)
    return shp


def _pic_full(slide, path):
    try:
        slide.shapes.add_picture(path, 0, 0, SLIDE_W, SLIDE_H)
        return True
    except Exception:
        return False


def _add_textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tb, tf


def _eyebrow(slide, text, color, l, t, w):
    """小号眉标（章节/栏目标签），营造层级。"""
    tb, tf = _add_textbox(slide, l, t, w, Inches(0.4), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    _set_font(r, "微软雅黑", 13, color, bold=True)
    return tb


def _footer(slide, style, page_no, source_refs=None, note=None):
    """页脚：左=来源标注，右=页码；细分割线。"""
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(7.05), SLIDE_W - Inches(1.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(style["accent_soft"])
    line.line.fill.background()
    line.shadow.inherit = False

    left = "　".join(f"来源：{s}" for s in (source_refs or [])) or (note or "")
    if left:
        tb, tf = _add_textbox(slide, Inches(0.7), Inches(7.06), SLIDE_W - Inches(3.0), Inches(0.35), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = left[:60]
        _set_font(r, "微软雅黑", 10, _rgb(style["body_color"]))

    if page_no:
        tb, tf = _add_textbox(slide, SLIDE_W - Inches(2.3), Inches(7.06), Inches(1.6), Inches(0.35), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{page_no:02d}"
        _set_font(r, "微软雅黑", 11, _rgb(style["accent"]), bold=True)


# ───────────────────────── 参数化渲染 ─────────────────────────

def _cover_bg(slide, style, cover_img):
    if cover_img and _pic_full(slide, cover_img):
        _rect_alpha(slide, 0, 0, SLIDE_W, SLIDE_H, style.get("scrim", "0A1B2B"), 0.55)
    else:
        _bg(slide, style["cover_bg"])


def _render_cover(slide, style, content: Dict, cover_img=None):
    mode = style.get("mode", "swiss")
    _cover_bg(slide, style, cover_img)
    title = content.get("title") or "毕业答辩"
    accent_c = _rgb(style["accent"])

    if mode == "swiss":
        # 强对比 / 左对齐 / 网格感（瑞士国际主义）
        tb, tf = _add_textbox(slide, Inches(0.9), Inches(1.9), SLIDE_W - Inches(1.8), Inches(0.5), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "GRADUATION DEFENSE"; _set_font(r, "微软雅黑", 14, accent_c, bold=True)
        tb2, tf2 = _add_textbox(slide, Inches(0.9), Inches(2.4), SLIDE_W - Inches(1.8), Inches(2.4), MSO_ANCHOR.TOP)
        p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = title; _set_font(r2, style["font_title"], 46, _WHITE, bold=True)
        rule = slide.shapes.add_shape(1, Inches(0.9), Inches(4.95), Inches(2.6), Inches(0.06))
        rule.fill.solid(); rule.fill.fore_color.rgb = accent_c; rule.line.fill.background(); rule.shadow.inherit = False
        tbs, tfs = _add_textbox(slide, Inches(0.9), Inches(5.15), SLIDE_W - Inches(1.8), Inches(0.6))
        ps = tfs.paragraphs[0]; ps.alignment = PP_ALIGN.LEFT
        rs = ps.add_run(); rs.text = "毕业答辩汇报"; _set_font(rs, style["font_body"], 20, RGBColor(0xEE, 0xEE, 0xEE))
        tbi, tfi = _add_textbox(slide, Inches(0.9), Inches(5.85), SLIDE_W - Inches(1.8), Inches(0.8))
        pi = tfi.paragraphs[0]; pi.alignment = PP_ALIGN.LEFT
        ri = pi.add_run(); ri.text = "专业：______   答辩人：______   导师：______"; _set_font(ri, style["font_body"], 15, RGBColor(0xDD, 0xDD, 0xDD))
    else:
        # 居中 / 衬线 / 编辑式（电子杂志）
        tb, tf = _add_textbox(slide, 0, Inches(2.0), SLIDE_W, Inches(0.5), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "毕 业 答 辩"; _set_font(r, "微软雅黑", 14, RGBColor(0xCF, 0xDA, 0xE8), bold=True)
        tb2, tf2 = _add_textbox(slide, Inches(1.0), Inches(2.6), SLIDE_W - Inches(2.0), Inches(2.0), MSO_ANCHOR.MIDDLE)
        p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = title; _set_font(r2, style["font_title"], 44, _WHITE, bold=True)
        rule = slide.shapes.add_shape(1, Inches(5.66), Inches(4.75), Inches(2.0), Inches(0.05))
        rule.fill.solid(); rule.fill.fore_color.rgb = accent_c; rule.line.fill.background(); rule.shadow.inherit = False
        tbs, tfs = _add_textbox(slide, Inches(1.0), Inches(4.9), SLIDE_W - Inches(2.0), Inches(0.6))
        ps = tfs.paragraphs[0]; ps.alignment = PP_ALIGN.CENTER
        rs = ps.add_run(); rs.text = "毕业答辩汇报"; _set_font(rs, style["font_body"], 20, RGBColor(0xEE, 0xEE, 0xEE))
        tbi, tfi = _add_textbox(slide, Inches(1.0), Inches(5.6), SLIDE_W - Inches(2.0), Inches(0.8))
        pi = tfi.paragraphs[0]; pi.alignment = PP_ALIGN.CENTER
        ri = pi.add_run(); ri.text = "专业：______   答辩人：______   导师：______"; _set_font(ri, style["font_body"], 15, RGBColor(0xDD, 0xDD, 0xDD))


def _render_section(slide, style, content: Dict, section_img=None, idx=None):
    mode = style.get("mode", "swiss")
    if section_img and _pic_full(slide, section_img):
        _rect_alpha(slide, 0, 0, SLIDE_W, SLIDE_H, style.get("scrim", "0A1B2B"), 0.45)
        title_color = _WHITE
    else:
        _bg(slide, style["bg_alt"])
        title_color = _rgb(style["accent"])
    # 大号章节序号（取该节在序列中的位置，由调用方传 idx）
    num = f"{idx:02d}" if idx else ""
    if num:
        if mode == "swiss":
            ntb, ntf = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(4.6), Inches(3.2), MSO_ANCHOR.MIDDLE)
            np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.LEFT
            nr = np.add_run(); nr.text = num; _set_font(nr, style["font_title"], 150, _rgb(style["accent_soft"]), bold=True)
        else:
            ntb, ntf = _add_textbox(slide, SLIDE_W - Inches(5.0), Inches(1.2), Inches(4.4), Inches(3.2), MSO_ANCHOR.MIDDLE)
            np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.RIGHT
            nr = np.add_run(); nr.text = num; _set_font(nr, style["font_title"], 150, _rgb(style["accent_soft"]), bold=True)
    _eyebrow(slide, "CHAPTER", _rgb(style["accent"]), Inches(0.9), Inches(4.7), SLIDE_W - Inches(1.8))
    rule = slide.shapes.add_shape(1, Inches(0.9), Inches(4.55), Inches(1.8), Inches(0.05))
    rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(style["accent"]); rule.line.fill.background(); rule.shadow.inherit = False
    tb, tf = _add_textbox(slide, Inches(0.9), Inches(5.05), SLIDE_W - Inches(1.8), Inches(1.6), MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = content.get("title", ""); _set_font(r, style["font_title"], 34, title_color, bold=True)


def _title_header(slide, style, title, page_idx=None):
    """内容页统一页眉：眉标+标题+强调横杠。"""
    _eyebrow(slide, "CONTENT", _rgb(style["accent"]), Inches(0.9), Inches(0.55), SLIDE_W - Inches(1.8))
    tb, tf = _add_textbox(slide, Inches(0.9), Inches(0.95), SLIDE_W - Inches(1.8), Inches(0.9), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    _set_font(r, style["font_title"], style["font_size_title"], _DARK, bold=True)
    bar = slide.shapes.add_shape(1, Inches(0.9), Inches(1.85), Inches(2.2), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(style["accent"]); bar.line.fill.background(); bar.shadow.inherit = False
    # 签名强调：左侧细竖条，统一视觉指纹
    sig = slide.shapes.add_shape(1, Inches(0.55), Inches(0.95), Inches(0.08), Inches(0.9))
    sig.fill.solid(); sig.fill.fore_color.rgb = _rgb(style["accent"]); sig.line.fill.background(); sig.shadow.inherit = False


def _bullets(tf, items: List[str], style, size=None, color=None):
    size = size or style["font_size_body"]
    color = color or _rgb(style["body_color"])
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(14)
        p.space_before = Pt(0)
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = "● "
        _set_font(r1, style["font_body"], size, _rgb(style["accent"]), bold=True)
        r2 = p.add_run(); r2.text = it
        _set_font(r2, style["font_body"], size, color)


def _render_bullets(slide, style, content: Dict, page_no=None):
    _bg(slide, style["bg"])
    _title_header(slide, style, content.get("title", ""))
    tb, tf = _add_textbox(slide, Inches(0.9), Inches(2.2), SLIDE_W - Inches(1.8), Inches(4.6))
    _bullets(tf, content.get("bullets", []), style)
    _footer(slide, style, page_no, content.get("source_refs"), content.get("note"))


def _render_two_column(slide, style, content: Dict, page_no=None):
    _bg(slide, style["bg"])
    _title_header(slide, style, content.get("title", ""))
    bullets = content.get("bullets", [])
    mid = (len(bullets) + 1) // 2
    # 左栏卡片
    _rect(slide, Inches(0.9), Inches(2.2), Inches(5.6), Inches(4.4), style["surface"], line_hex=style["accent_soft"], line_w=1)
    lb, ltf = _add_textbox(slide, Inches(1.2), Inches(2.5), Inches(5.0), Inches(3.8))
    _bullets(ltf, bullets[:mid] or ["—"], style)
    _rect(slide, Inches(6.83), Inches(2.2), Inches(5.6), Inches(4.4), style["surface"], line_hex=style["accent_soft"], line_w=1)
    rb, rtf = _add_textbox(slide, Inches(7.13), Inches(2.5), Inches(5.0), Inches(3.8))
    _bullets(rtf, bullets[mid:] or ["—"], style)
    _footer(slide, style, page_no, content.get("source_refs"), content.get("note"))


def _render_image_right(slide, style, content: Dict, image: Optional[str], page_no=None):
    _bg(slide, style["bg"])
    _title_header(slide, style, content.get("title", ""))
    tb, tf = _add_textbox(slide, Inches(0.9), Inches(2.2), Inches(6.3), Inches(4.4))
    _bullets(tf, content.get("bullets", []), style)
    rx = Inches(7.5)
    rw = Inches(5.0)
    rh = Inches(4.4)
    # 白色卡片衬底
    _rect(slide, rx - Inches(0.15), Inches(2.05), rw + Inches(0.3), rh + Inches(0.3), style["surface"], line_hex=style["accent_soft"], line_w=1)
    if image:
        try:
            slide.shapes.add_picture(image, rx, Inches(2.2), rw, rh)
        except Exception:
            _placeholder(slide, style, rx, Inches(2.2), rw, rh)
    else:
        _placeholder(slide, style, rx, Inches(2.2), rw, rh)
    _footer(slide, style, page_no, content.get("source_refs"), content.get("note"))


def _placeholder(slide, style, l, t, w, h):
    ph = slide.shapes.add_shape(1, l, t, w, h)
    ph.fill.solid(); ph.fill.fore_color.rgb = _rgb(style["accent_soft"])
    ph.line.color.rgb = _rgb(style["accent"]); ph.shadow.inherit = False
    ptf = ph.text_frame; ptf.word_wrap = True
    pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    pr = pp.add_run(); pr.text = "图片占位\n（将素材图放入 assets 目录可自动填入）"
    _set_font(pr, style["font_body"], 14, _rgb(style["accent"]))


def _render_closing(slide, style, content: Dict):
    _bg(slide, style["cover_bg"])
    tb, tf = _add_textbox(slide, Inches(1.0), Inches(3.0), SLIDE_W - Inches(2.0), Inches(1.6), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = content.get("title", "感谢聆听，请批评指正")
    _set_font(r, style["font_title"], 38, _WHITE, bold=True)


# ───────────────────────── 真实母版渲染 ─────────────────────────

def _open_master_prs(master_path: str, out_path: str):
    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    shutil.copy(master_path, out_path)
    prs = Presentation(out_path)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)
    return prs


def _render_master(slide, style, s, imgs, img_state, page_no=None, section_idx=None):
    layout = s.get("layout", "bullets")
    title = s.get("title", "")
    if layout == "cover":
        _render_cover(slide, style, s, imgs["cover"])
        return
    if layout == "closing":
        _render_closing(slide, style, s)
        return
    if layout == "section":
        _render_section(slide, style, s, imgs["section"], idx=section_idx)
        return

    # 内容页：在母版主题上铺白色卡片，保证可读
    _rect(slide, Inches(0.45), Inches(0.45), SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.9), style["surface"])
    _title_header(slide, style, title)
    bullets = s.get("bullets", [])
    if layout == "two_column":
        mid = (len(bullets) + 1) // 2
        lb, ltf = _add_textbox(slide, Inches(0.9), Inches(2.2), Inches(5.6), Inches(4.4))
        _bullets(ltf, bullets[:mid] or ["—"], style)
        rb, rtf = _add_textbox(slide, Inches(6.9), Inches(2.2), Inches(5.6), Inches(4.4))
        _bullets(rtf, bullets[mid:] or ["—"], style)
    elif layout == "image_right":
        tb2, tf2 = _add_textbox(slide, Inches(0.9), Inches(2.2), Inches(6.3), Inches(4.4))
        _bullets(tf2, bullets, style)
        img = imgs["content"][img_state[0]] if img_state[0] < len(imgs["content"]) else None
        rx, rw, rh = Inches(7.5), Inches(5.0), Inches(4.4)
        _rect(slide, rx - Inches(0.15), Inches(2.05), rw + Inches(0.3), rh + Inches(0.3), style["surface"], line_hex=style["accent_soft"], line_w=1)
        if img:
            try:
                slide.shapes.add_picture(img, rx, Inches(2.2), rw, rh)
                img_state[0] += 1
            except Exception:
                _placeholder(slide, style, rx, Inches(2.2), rw, rh)
        else:
            _placeholder(slide, style, rx, Inches(2.2), rw, rh)
    else:
        tb2, tf2 = _add_textbox(slide, Inches(0.9), Inches(2.2), SLIDE_W - Inches(1.8), Inches(4.4))
        _bullets(tf2, bullets, style)
    _footer(slide, style, page_no, s.get("source_refs"), s.get("note"))


# ───────────────────────── 入口 ─────────────────────────

def build_pptx(
    content: Dict,
    style_id: str,
    out_path: str,
    assets_dir: Optional[str] = None,
) -> str:
    styles = _load_styles()
    style = styles.get(style_id, styles["minimal_academic"])
    master = _real_master(content.get("discipline", ""), style_id)
    imgs = _resolve_images(assets_dir)

    if master and os.path.exists(master):
        prs = _open_master_prs(master, out_path)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    slides_data = content.get("slides", [])
    blank = _blank_layout(prs)
    page_counter = 0
    section_counter = 0

    for s in slides_data:
        slide = prs.slides.add_slide(blank)
        layout = s.get("layout", "bullets")
        if master:
            img_state = [0]
            if layout == "section":
                section_counter += 1
                _render_master(slide, style, s, imgs, img_state, page_no=None, section_idx=section_counter)
            else:
                _render_master(slide, style, s, imgs, img_state, page_no=None)
            continue
        if layout == "cover":
            _render_cover(slide, style, s, imgs["cover"])
        elif layout == "section":
            section_counter += 1
            _render_section(slide, style, s, imgs["section"], idx=section_counter)
        elif layout == "two_column":
            page_counter += 1
            _render_two_column(slide, style, s, page_no=page_counter)
        elif layout == "image_right":
            page_counter += 1
            img = imgs["content"][0] if imgs["content"] else None
            if img:
                imgs["content"] = imgs["content"][1:]
            _render_image_right(slide, style, s, img, page_no=page_counter)
        elif layout == "closing":
            _render_closing(slide, style, s)
        else:
            page_counter += 1
            _render_bullets(slide, style, s, page_no=page_counter)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    return out_path


def _blank_layout(prs):
    for lay in prs.slide_masters[0].slide_layouts:
        if len(lay.placeholders) == 0:
            return lay
    lays = prs.slide_masters[0].slide_layouts
    return lays[6] if len(lays) > 6 else lays[0]
